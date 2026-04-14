"""
Rapport J+8 - The Pulse  (PowerPoint version)
Génère un fichier .pptx éditable
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

IMG_DIR = os.path.join(os.path.dirname(__file__), "static", "images")
LOGO_DARK = os.path.join(IMG_DIR, "thepulse-noir.png")
LOGO_WHITE = os.path.join(IMG_DIR, "thepulse-blanc.png")
OUTPUT = os.path.join(os.path.dirname(__file__), "Rapport_J7_ThePulse_Partenaires.pptx")

# Colors
ACCENT = RGBColor(0x00, 0xd4, 0xaa)
DARK = RGBColor(0x0f, 0x17, 0x2a)
MUTED = RGBColor(0x64, 0x74, 0x8b)
ORANGE = RGBColor(0xf9, 0x73, 0x16)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x22, 0xc5, 0x5e)
WHITE_C = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xf8, 0xfa, 0xfc)
CARD_BG = RGBColor(0xf1, 0xf5, 0xf9)
SLATE = RGBColor(0x47, 0x55, 0x69)

PARTNER_LOGOS = [
    (os.path.join(IMG_DIR, "um6p_logo.png"), "UM6P"),
    (os.path.join(IMG_DIR, "logo_omtpme-13.png"), "OMTPME"),
    (os.path.join(IMG_DIR, "tamwilcom_logo.png"), "Tamwilcom"),
    (os.path.join(IMG_DIR, "amic_logo.png"), "AMIC"),
    (os.path.join(IMG_DIR, "MTN LOGO.svg"), "MTN"),
]


def add_bg(slide, color):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, corner_radius=0):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if corner_radius:
        try:
            shape.adjustments[0] = corner_radius
        except Exception:
            pass
    return shape


def add_text(slide, left, top, width, height, text, font_size=12, bold=False, color=DARK, alignment=PP_ALIGN.LEFT):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def add_kpi_card(slide, left, top, width, height, value, label, accent_color=ACCENT):
    """Add a KPI card with value and label."""
    # Card background
    card = add_rect(slide, left, top, width, height, CARD_BG, corner_radius=0.08)
    # Accent bar
    add_rect(slide, left, top, width, Pt(4), accent_color)
    # Value
    add_text(slide, left, top + Pt(12), width, Pt(30), str(value),
             font_size=24, bold=True, color=DARK, alignment=PP_ALIGN.CENTER)
    # Label
    add_text(slide, left, top + height - Pt(22), width, Pt(18), label,
             font_size=8, bold=False, color=MUTED, alignment=PP_ALIGN.CENTER)


def add_table(slide, left, top, headers, rows, col_widths_inches):
    """Add a styled table."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top,
                                          sum(Inches(w) for w in col_widths_inches), Inches(0.3 * n_rows))
    table = table_shape.table

    # Set column widths
    for i, w in enumerate(col_widths_inches):
        table.columns[i].width = Inches(w)

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = WHITE_C
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8)
                p.font.bold = (ci == 0)
                p.font.color.rgb = DARK
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE_C

    return table_shape


def add_page_header(slide, title, right_text):
    """Add standard page header with dark bar and logo."""
    add_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.55), DARK)
    add_rect(slide, Inches(0), Inches(0.55), Inches(10), Pt(4), ACCENT)
    # Logo
    try:
        slide.shapes.add_picture(LOGO_DARK, Inches(0.3), Inches(0.1), height=Inches(0.3))
    except Exception:
        pass
    add_text(slide, Inches(1.1), Inches(0.12), Inches(6), Inches(0.4), title,
             font_size=14, bold=True, color=WHITE_C)
    add_text(slide, Inches(7), Inches(0.12), Inches(2.7), Inches(0.4), right_text,
             font_size=9, color=RGBColor(0x94, 0xa3, 0xb8), alignment=PP_ALIGN.RIGHT)


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 1 - COVER
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_bg(slide, DARK)

    # Accent top bar
    add_rect(slide, Inches(0), Inches(0), Inches(10), Pt(6), ACCENT)

    # Logo
    try:
        slide.shapes.add_picture(LOGO_WHITE, Inches(3.2), Inches(0.4), width=Inches(3.6))
    except Exception:
        pass

    # Title
    add_text(slide, Inches(0.5), Inches(1.8), Inches(9), Inches(0.7), "Rapport J+8",
             font_size=40, bold=True, color=WHITE_C)
    add_text(slide, Inches(0.5), Inches(2.5), Inches(9), Inches(0.6), "Bilan de Lancement",
             font_size=30, bold=True, color=WHITE_C)

    # Accent underline
    add_rect(slide, Inches(0.5), Inches(3.15), Inches(1.2), Pt(4), ACCENT)

    # Subtitle
    add_text(slide, Inches(0.5), Inches(3.4), Inches(8), Inches(0.4),
             "Plateforme de données de l\u2019écosystème startup marocain",
             font_size=14, color=RGBColor(0x94, 0xa3, 0xb8))

    # Date
    add_text(slide, Inches(0.5), Inches(4.1), Inches(5), Inches(0.3),
             "Période : 6 - 13 avril 2026", font_size=13, bold=True, color=ACCENT)
    add_text(slide, Inches(0.5), Inches(4.4), Inches(5), Inches(0.3),
             "8 jours après le lancement officiel", font_size=11, color=RGBColor(0x94, 0xa3, 0xb8))

    # Partner logos label
    add_text(slide, Inches(0.5), Inches(5.1), Inches(3), Inches(0.3),
             "Rapport partenaires", font_size=10, bold=True, color=SLATE)

    # Partner logos
    x_start = 0.5
    for i, (logo_path, name) in enumerate(PARTNER_LOGOS):
        lx = x_start + i * 1.8
        try:
            slide.shapes.add_picture(logo_path, Inches(lx), Inches(5.5), height=Inches(0.55))
        except Exception:
            pass
        add_text(slide, Inches(lx), Inches(6.1), Inches(1.5), Inches(0.2),
                 name, font_size=7, color=MUTED)

    # MESC box
    mesc_shape = add_rect(slide, Inches(x_start + len(PARTNER_LOGOS) * 1.8), Inches(5.5),
                          Inches(0.55), Inches(0.55), PURPLE, corner_radius=0.1)
    add_text(slide, Inches(x_start + len(PARTNER_LOGOS) * 1.8), Inches(5.55),
             Inches(0.55), Inches(0.45), "M", font_size=18, bold=True, color=WHITE_C, alignment=PP_ALIGN.CENTER)
    add_text(slide, Inches(x_start + len(PARTNER_LOGOS) * 1.8 - 0.1), Inches(6.1),
             Inches(0.8), Inches(0.2), "MESC", font_size=7, color=MUTED)

    # Footer
    add_text(slide, Inches(0.5), Inches(6.8), Inches(5), Inches(0.2),
             "The Pulse  |  UM6P  |  thepulse.ma", font_size=9, color=SLATE)
    add_text(slide, Inches(0.5), Inches(7.0), Inches(5), Inches(0.2),
             "Document confidentiel - Rapport partenaires", font_size=9, color=SLATE)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 2 - VUE D'ENSEMBLE
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_page_header(slide, "1.  Vue d\u2019ensemble de la plateforme", "The Pulse  |  J+8")

    add_text(slide, Inches(0.4), Inches(0.75), Inches(6), Inches(0.3),
             "Base de données écosystème", font_size=13, bold=True, color=DARK)
    add_text(slide, Inches(0.4), Inches(1.0), Inches(8), Inches(0.25),
             "Données collectées, structurées et mises en ligne sur thepulse.ma",
             font_size=9, color=MUTED)

    # KPI Row 1
    kpi_w = Inches(2.2)
    kpi_h = Inches(0.75)
    kpi_gap = Inches(0.15)
    kpi_y = Inches(1.35)
    kpis_row1 = [
        ("2 014", "STARTUPS", ACCENT),
        ("1 333", "FONDATEURS", PURPLE),
        ("51", "INVESTISSEURS", BLUE),
        ("$226.8M", "FONDS LEVÉS", ORANGE),
    ]
    for i, (val, lbl, clr) in enumerate(kpis_row1):
        add_kpi_card(slide, Inches(0.4) + i * (kpi_w + kpi_gap), kpi_y, kpi_w, kpi_h, val, lbl, clr)

    # KPI Row 2
    kpi_y2 = Inches(2.25)
    kpis_row2 = [
        ("171", "TOURS DE TABLE", GREEN),
        ("45", "INCUBATEURS", BLUE),
        ("1 265", "LIENS STARTUP-FONDATEUR", PURPLE),
        ("21", "RESSOURCES", MUTED),
    ]
    for i, (val, lbl, clr) in enumerate(kpis_row2):
        add_kpi_card(slide, Inches(0.4) + i * (kpi_w + kpi_gap), kpi_y2, kpi_w, kpi_h, val, lbl, clr)

    # Detail table
    add_text(slide, Inches(0.4), Inches(3.2), Inches(6), Inches(0.3),
             "Détail de la base de données", font_size=13, bold=True, color=DARK)

    headers = ["CATÉGORIE", "NOMBRE", "DÉTAIL"]
    rows = [
        ["Startups", "2 014", "Répertoire complet des startups marocaines"],
        ["Fondateurs", "1 333", "948 startups avec fondateur(s) identifié(s)"],
        ["Investisseurs", "51", "Fonds d\u2019investissement actifs au Maroc"],
        ["Incubateurs / Accélérateurs", "45", "Programmes d\u2019accompagnement"],
        ["Tours de financement", "171", "Levées de fonds documentées ($226.8M total)"],
        ["Experts référencés", "6", "Consultants et mentors"],
        ["Ressources", "21", "Guides, rapports, outils"],
        ["Articles / Actualités", "9", "Contenu éditorial"],
        ["Membres inscrits", "44", "Comptes créés sur la plateforme (entrepreneurs, investisseurs, experts)"],
    ]
    add_table(slide, Inches(0.4), Inches(3.55), headers, rows, [2.2, 0.9, 6.1])

    add_text(slide, Inches(0.4), Inches(7.1), Inches(9), Inches(0.2),
             "* Enrichissement continu : scraping multi-sources + validation manuelle. Objectif : 100% des startups avec fondateur(s) identifié(s).",
             font_size=8, color=MUTED)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 3 - TRAFIC & AUDIENCE
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_page_header(slide, "2.  Trafic & Audience  (Google Analytics)", "6 - 13 avril 2026")

    add_text(slide, Inches(0.4), Inches(0.75), Inches(6), Inches(0.3),
             "Indicateurs clés de trafic  (7 jours)", font_size=13, bold=True, color=DARK)

    # KPI Row 1
    kpi_w3 = Inches(3.0)
    kpi_y = Inches(1.1)
    kpis = [
        ("2 665", "SESSIONS", ACCENT),
        ("1 855", "UTILISATEURS ACTIFS", PURPLE),
        ("1 748", "NOUVEAUX UTILISATEURS", BLUE),
    ]
    for i, (val, lbl, clr) in enumerate(kpis):
        add_kpi_card(slide, Inches(0.4) + i * (kpi_w3 + kpi_gap), kpi_y, kpi_w3, kpi_h, val, lbl, clr)

    kpi_y2 = Inches(2.0)
    kpis2 = [
        ("57,7%", "TAUX D\u2019ENGAGEMENT", GREEN),
        ("1m 34s", "DURÉE MOY. SESSION", ORANGE),
        ("30 683", "ÉVÉNEMENTS", MUTED),
    ]
    for i, (val, lbl, clr) in enumerate(kpis2):
        add_kpi_card(slide, Inches(0.4) + i * (kpi_w3 + kpi_gap), kpi_y2, kpi_w3, kpi_h, val, lbl, clr)

    # Sources d'acquisition
    add_text(slide, Inches(0.4), Inches(2.95), Inches(6), Inches(0.3),
             "Sources d\u2019acquisition", font_size=13, bold=True, color=DARK)

    channels = [
        ("Direct", "1 378 (51,7%)"),
        ("Recherche organique (Google)", "1 063 (39,9%)"),
        ("Réseaux sociaux (LinkedIn)", "114 (4,3%)"),
        ("Referrals (Medias24, etc.)", "63 (2,4%)"),
        ("Autres", "47 (1,8%)"),
    ]
    for i, (label, val) in enumerate(channels):
        y_pos = Inches(3.3) + i * Inches(0.22)
        add_text(slide, Inches(0.5), y_pos, Inches(3.5), Inches(0.2),
                 f"{label}  —  {val}", font_size=9, color=DARK)

    # Géographie table
    add_text(slide, Inches(0.4), Inches(4.55), Inches(6), Inches(0.3),
             "Répartition géographique des utilisateurs", font_size=13, bold=True, color=DARK)

    headers = ["PAYS", "UTILISATEURS", "PART", "TAUX ENGAGEMENT", "DURÉE MOY."]
    rows = [
        ["Maroc", "1 440", "77,6%", "55,2%", "2m 18s"],
        ["France", "151", "8,1%", "61,5%", "1m 42s"],
        ["États-Unis", "56", "3,0%", "35,8%", "36s"],
        ["Canada", "33", "1,8%", "59,4%", "1m 04s"],
        ["Pays-Bas", "19", "1,0%", "76,1%", "3m 10s"],
        ["Royaume-Uni", "14", "0,8%", "56,3%", "33s"],
        ["Allemagne", "11", "0,6%", "60,0%", "1m 51s"],
        ["Espagne", "3", "0,2%", "50,0%", "45s"],
    ]
    add_table(slide, Inches(0.4), Inches(4.9), headers, rows, [1.6, 1.1, 0.8, 1.3, 1.1])

    add_text(slide, Inches(0.4), Inches(7.2), Inches(9), Inches(0.2),
             "Source : Google Analytics (G-9TTHWMF6L0) - Période du 6 au 13 avril 2026",
             font_size=8, color=MUTED)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 4 - PAGES & ENGAGEMENT
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_page_header(slide, "3.  Pages les plus visitées & Engagement", "6 - 13 avril 2026")

    add_text(slide, Inches(0.4), Inches(0.75), Inches(6), Inches(0.3),
             "Pages les plus consultées", font_size=13, bold=True, color=DARK)

    headers = ["PAGE", "VUES", "UTILISATEURS", "REBOND", "VS MOYENNE*"]
    rows = [
        ["Répertoire Startups", "3 067", "690", "18,5%", "\u25b2 Page #1"],
        ["Page d\u2019accueil", "1 892", "1 310", "35,2%", "Rebond moyen"],
        ["Rejoindre The Pulse", "1 010", "725", "17,9%", "\u25b2 Fort intérêt"],
        ["Répertoire Fondateurs", "576", "215", "11,0%", "\u25b2 Très engagé"],
        ["Répertoire Investisseurs", "403", "176", "10,1%", "\u25b2 Très engagé"],
        ["Se connecter", "302", "240", "12,8%", "\u25b2 Engagé"],
        ["Co-Fondateurs", "282", "195", "9,2%", "\u25b2\u25b2 Excellent"],
    ]
    add_table(slide, Inches(0.4), Inches(1.1), headers, rows, [2.4, 0.8, 1.1, 1.0, 1.5])

    add_text(slide, Inches(0.4), Inches(3.55), Inches(9), Inches(0.2),
             "* Moyenne taux de rebond SaaS B2B : 40-60% (source : CXL / Contentsquare 2025). The Pulse est largement en dessous.",
             font_size=7, color=MUTED)

    # Tendance des visites (J1→J7)
    add_text(slide, Inches(0.4), Inches(3.85), Inches(6), Inches(0.3),
             "Tendance des visites (J1 \u2192 J8)", font_size=13, bold=True, color=DARK)

    trend_days = [
        ("J1 (6 avr)", "180", MUTED),
        ("J2 (7 avr)", "220", MUTED),
        ("J3 (8 avr)", "290", BLUE),
        ("J4 (9 avr)", "520", ACCENT),
        ("J5 (10 avr)", "480", ACCENT),
        ("J6 (11 avr)", "410", BLUE),
        ("J7 (12 avr)", "341", BLUE),
        ("J8 (13 avr)", "224", MUTED),
    ]
    bar_max = 520
    bar_x = Inches(0.5)
    bar_w_max = Inches(5.0)
    for i, (day, sessions, clr) in enumerate(trend_days):
        y_pos = Inches(4.2) + i * Inches(0.32)
        # Day label
        add_text(slide, bar_x, y_pos, Inches(1.3), Inches(0.25),
                 day, font_size=8, color=DARK)
        # Bar
        ratio = int(sessions.replace(" ", "")) / bar_max
        bw = max(Inches(0.2), bar_w_max * ratio)
        add_rect(slide, bar_x + Inches(1.35), y_pos + Inches(0.02), int(bw), Inches(0.2), clr, corner_radius=0.15)
        # Value
        add_text(slide, bar_x + Inches(1.35) + int(bw) + Inches(0.1), y_pos, Inches(0.8), Inches(0.25),
                 f"{sessions} sess.", font_size=7, bold=True, color=DARK)

    add_text(slide, Inches(0.4), Inches(6.5), Inches(9), Inches(0.2),
             "\u25b2 Pic J4 (9 avril) : +189% vs J1 \u2014 effet du lancement LinkedIn + partage communauté",
             font_size=8, bold=True, color=ACCENT)

    # Insight box
    insight_shape = add_rect(slide, Inches(0.4), Inches(6.8), Inches(9.2), Inches(0.55),
                              RGBColor(0xec, 0xfd, 0xf5), corner_radius=0.05)
    add_rect(slide, Inches(0.4), Inches(6.8), Pt(4), Inches(0.55), ACCENT)
    add_text(slide, Inches(0.6), Inches(6.82), Inches(2), Inches(0.2),
             "Point clé", font_size=10, bold=True, color=DARK)
    add_text(slide, Inches(0.6), Inches(7.02), Inches(8.8), Inches(0.35),
             "Taux de rebond moyen de 16,6% vs 40-60% pour un SaaS B2B \u2014 les visiteurs explorent activement la plateforme. 44 inscriptions en 8 jours avec un pic à J4.",
             font_size=8.5, color=RGBColor(0x33, 0x41, 0x55))

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 5 - COMMUNAUTÉ & INSCRITS
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_page_header(slide, "4.  Communauté & Membres inscrits", "Au 13 avril 2026")

    add_text(slide, Inches(0.4), Inches(0.75), Inches(6), Inches(0.3),
             "Inscriptions à la plateforme", font_size=13, bold=True, color=DARK)

    kpi_y = Inches(1.1)
    kpis = [
        ("48", "INSCRITS TOTAL", ACCENT),
        ("22", "CONFIRMÉS", GREEN),
        ("26", "EN ATTENTE", ORANGE),
        ("11", "AVEC PHOTO", PURPLE),
    ]
    for i, (val, lbl, clr) in enumerate(kpis):
        add_kpi_card(slide, Inches(0.4) + i * (kpi_w + kpi_gap), kpi_y, kpi_w, kpi_h, val, lbl, clr)

    # Répartition par rôle
    add_text(slide, Inches(0.4), Inches(2.05), Inches(6), Inches(0.3),
             "Répartition par rôle", font_size=13, bold=True, color=DARK)

    roles = [
        ("Entrepreneurs", "35 (72,9%)"),
        ("Investisseurs", "5 (10,4%)"),
        ("Experts", "5 (10,4%)"),
        ("Programmes / Incubateurs", "3 (6,3%)"),
    ]
    for i, (label, val) in enumerate(roles):
        y_pos = Inches(2.4) + i * Inches(0.22)
        add_text(slide, Inches(0.5), y_pos, Inches(5), Inches(0.2),
                 f"{label}  —  {val}", font_size=10, color=DARK)

    # Inscriptions par jour
    add_text(slide, Inches(0.4), Inches(3.4), Inches(6), Inches(0.3),
             "Inscriptions par jour", font_size=13, bold=True, color=DARK)

    days = [
        ("6 avril (lancement)", "0"),
        ("7 avril", "0"),
        ("8 avril", "0"),
        ("9 avril", "29"),
        ("10 avril", "8"),
        ("11 avril", "4"),
        ("12 avril", "3"),
        ("13 avril", "4"),
    ]
    for i, (day, count) in enumerate(days):
        y_pos = Inches(3.75) + i * Inches(0.2)
        val_str = count if count != "0" else "—"
        add_text(slide, Inches(0.5), y_pos, Inches(4), Inches(0.2),
                 f"{day}  :  {val_str}", font_size=9, color=DARK)

    # Engagement communautaire
    add_text(slide, Inches(0.4), Inches(5.3), Inches(6), Inches(0.3),
             "Engagement communautaire", font_size=13, bold=True, color=DARK)

    headers = ["INDICATEUR", "VALEUR", "COMMENTAIRE"]
    rows = [
        ["Messages directs échangés", "10", "Messagerie interne entre membres"],
        ["Publications newsfeed", "14", "Posts de la communauté"],
        ["Projets co-fondateur", "6", "Recherches de co-fondateurs actives"],
        ["Membres avec mot de passe", "8", "Comptes pleinement activés"],
    ]
    add_table(slide, Inches(0.4), Inches(5.65), headers, rows, [2.5, 0.9, 5.0])

    # CTA box for partners
    cta_shape = add_rect(slide, Inches(0.4), Inches(6.85), Inches(9.2), Inches(0.55),
                          PURPLE, corner_radius=0.06)
    add_text(slide, Inches(0.6), Inches(6.88), Inches(8.8), Inches(0.25),
             "\U0001f91d  Appel aux partenaires", font_size=11, bold=True, color=WHITE_C)
    add_text(slide, Inches(0.6), Inches(7.12), Inches(8.8), Inches(0.25),
             "Aidez-nous \u00e0 promouvoir The Pulse ! Partagez la plateforme avec vos r\u00e9seaux, startups et communaut\u00e9s pour acc\u00e9l\u00e9rer l\u2019adoption.",
             font_size=9, color=RGBColor(0xe8, 0xda, 0xff))

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 6 - FONCTIONNALITÉS DÉPLOYÉES
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_page_header(slide, "5.  Fonctionnalités déployées depuis le lancement", "Produit & Communauté")

    add_text(slide, Inches(0.4), Inches(0.75), Inches(9), Inches(0.3),
             "Évolution de la plateforme en 8 jours", font_size=13, bold=True, color=DARK)
    add_text(slide, Inches(0.4), Inches(1.0), Inches(9), Inches(0.25),
             "Au-delà de la base de données, The Pulse est devenue une plateforme communautaire active.",
             font_size=9, color=MUTED)

    features = [
        (u"Comptes & Profils membres", PURPLE, [
            (u"Création de compte", u"Inscription via formulaires dédiés (entrepreneur, investisseur, expert)."),
            (u"Profils enrichis", u"Photo, bio, LinkedIn, compétences, parcours."),
            (u"Édition & mot de passe", u"Modification des infos et gestion du mot de passe."),
            (u"Badge « New Pulser »", u"Badge visuel sur les nouveaux membres inscrits."),
        ]),
        (u"Interactions & Communauté", ACCENT, [
            (u"Envoi de Pulses", u"Bouton « Envoyer un Pulse » pour saluer un membre."),
            (u"Messagerie instantanée (Inbox)", u"Conversations 1-to-1 entre membres."),
            (u"Newsfeed communautaire", u"Publications, annonces, opportunités."),
            (u"Projets co-fondateur", u"Espace dédié pour publier une recherche de co-fondateur."),
        ]),
        (u"Répertoires & Découverte", BLUE, [
            (u"Répertoire Startups", u"2 014 startups avec filtres secteur, ville, stade."),
            (u"Répertoire Fondateurs", u"1 333 fondateurs avec liens startups et LinkedIn."),
            (u"Répertoire Investisseurs", u"51 fonds avec thèses et tickets documentés."),
            (u"Talent Marketplace", u"Mise en relation talents / startups."),
        ]),
        (u"Administration & Opérations", ORANGE, [
            (u"Admin Pulsers", u"Panneau d\u2019administration des membres."),
            (u"Confirmation de compte", u"Workflow de validation des inscriptions."),
            (u"Gestion des inscriptions", u"Suivi des demandes en attente."),
            (u"Analytics & Rapports", u"Dashboards internes et génération de rapports PDF."),
        ]),
    ]

    col_w = Inches(4.6)
    col_h = Inches(2.7)
    positions = [
        (Inches(0.4), Inches(1.4)),
        (Inches(5.1), Inches(1.4)),
        (Inches(0.4), Inches(4.3)),
        (Inches(5.1), Inches(4.3)),
    ]
    for (cat_title, cat_color, items), (cx, cy) in zip(features, positions):
        # Category card
        add_rect(slide, cx, cy, col_w, col_h, WHITE_C, corner_radius=0.05)
        add_rect(slide, cx, cy, Inches(0.08), col_h, cat_color)
        add_text(slide, cx + Inches(0.2), cy + Inches(0.08), col_w - Inches(0.3), Inches(0.3),
                 cat_title, font_size=11, bold=True, color=DARK)
        for j, (name, desc) in enumerate(items):
            y_off = cy + Inches(0.45) + j * Inches(0.55)
            add_text(slide, cx + Inches(0.2), y_off, col_w - Inches(0.3), Inches(0.25),
                     u"\u2022 " + name, font_size=9, bold=True, color=DARK)
            add_text(slide, cx + Inches(0.35), y_off + Inches(0.22), col_w - Inches(0.4), Inches(0.3),
                     desc, font_size=7.5, color=SLATE)

    add_text(slide, Inches(0.4), Inches(7.15), Inches(9.2), Inches(0.2),
             u"The Pulse \u2014 d\u2019une base de donn\u00e9es \u00e0 une v\u00e9ritable plateforme communautaire en 8 jours.",
             font_size=8, color=MUTED, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 7 - SYNTHÈSE & PROCHAINES ÉTAPES
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_page_header(slide, "6.  Synthèse & Prochaines étapes", "Perspectives J+30")

    add_text(slide, Inches(0.4), Inches(0.75), Inches(6), Inches(0.3),
             "Synthèse J+8", font_size=13, bold=True, color=DARK)

    highlights = [
        ("Audience forte", "1 855 utilisateurs actifs et 2 665 sessions en 8 jours, principalement depuis le Maroc (78%)."),
        ("Acquisition organique", "39,6% du trafic provient de la recherche Google, signe d\u2019un bon référencement naturel."),
        ("Engagement élevé", "57,7% de taux d\u2019engagement, durée moyenne de session de 1m 34s, 48 inscriptions effectives en 8 jours."),
        ("Communauté naissante", "48 membres inscrits dont 35 entrepreneurs, 5 investisseurs, 5 experts et 3 programmes."),
        ("Base de données riche", "2 014 startups, 1 333 fondateurs, 51 investisseurs, 171 tours de financement documentés."),
        ("Rayonnement international", "Visiteurs de 8+ pays, avec une diaspora active (France 8,5%, USA 3,2%, Canada 1,8%)."),
    ]
    for i, (title, desc) in enumerate(highlights):
        y_pos = Inches(1.1) + i * Inches(0.4)
        # Bullet title + description
        txBox = slide.shapes.add_textbox(Inches(0.5), y_pos, Inches(9), Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = f"● {title}  "
        run1.font.size = Pt(9.5)
        run1.font.bold = True
        run1.font.color.rgb = DARK
        run2 = p.add_run()
        run2.text = desc
        run2.font.size = Pt(8.5)
        run2.font.color.rgb = SLATE

    # Axes d'amélioration
    add_text(slide, Inches(0.4), Inches(3.6), Inches(6), Inches(0.3),
             "Axes d\u2019amélioration identifiés", font_size=13, bold=True, color=DARK)

    improvements = [
        "Activation des 26 comptes en attente de confirmation (54% des inscrits)",
        "Enrichissement fondateurs : 1 068 startups sans fondateur identifié (objectif : -50% à J+30)",
        "Augmentation du contenu éditorial et des publications communautaires",
        "Déploiement d\u2019une stratégie d\u2019acquisition sur LinkedIn et les réseaux sociaux",
    ]
    for i, imp in enumerate(improvements):
        y_pos = Inches(3.95) + i * Inches(0.22)
        add_text(slide, Inches(0.5), y_pos, Inches(9), Inches(0.2),
                 f"● {imp}", font_size=9, color=SLATE)

    # Objectifs J+30
    add_text(slide, Inches(0.4), Inches(4.95), Inches(6), Inches(0.3),
             "Objectifs J+30", font_size=13, bold=True, color=DARK)

    headers = ["OBJECTIF", "ACTUEL", "CIBLE J+30", "PROGRESSION"]
    rows = [
        ["Membres inscrits", "48", "200", "24%"],
        ["Startups avec fondateur", "948", "1 500", "63%"],
        ["Publications communauté", "14", "50", "28%"],
        ["Sessions / semaine", "2 665", "5 000", "53%"],
    ]
    add_table(slide, Inches(0.4), Inches(5.3), headers, rows, [2.5, 1.0, 1.0, 1.5])

    # ══════════════════════════════════════════════════════════════════
    # SLIDE 8 - MENTORS & ASSIGNATIONS STARTUPS
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, LIGHT_BG)
    add_page_header(slide, "7.  Accompagnement : Mentors & Startups assign\u00e9es", "Programme de mentorat")

    add_text(slide, Inches(0.4), Inches(0.75), Inches(9), Inches(0.3),
             "Assignation Mentors \u2014 Startups", font_size=13, bold=True, color=DARK)
    add_text(slide, Inches(0.4), Inches(1.05), Inches(9), Inches(0.25),
             "Chaque mentor accompagne un portefeuille de startups s\u00e9lectionn\u00e9es pour un suivi personnalis\u00e9.",
             font_size=9, color=MUTED)

    mentors = [
        ("Ali El Amrani", PURPLE, [
            "IndusTwin",
            "IrrigSense",
            "AeroMind Morocco",
            "WAHGO Foods",
        ]),
        ("Youssef Mamou", ACCENT, [
            "Investing For Everyone",
            "Igudar",
            "XPredictia",
            "SmartDiag",
            "SyanaTek",
        ]),
        ("Tarik Fadli", BLUE, [
            "UM6P Instruments",
            "GeoHeritage",
            "G-ReLib",
        ]),
        ("Mohammed Damiri", ORANGE, [
            "DECAP",
            "GreenFlow",
        ]),
    ]

    card_w = Inches(4.5)
    card_h = Inches(2.5)
    mentor_positions = [
        (Inches(0.4), Inches(1.45)),
        (Inches(5.1), Inches(1.45)),
        (Inches(0.4), Inches(4.15)),
        (Inches(5.1), Inches(4.15)),
    ]
    for (mentor_name, mentor_color, startups), (mx, my) in zip(mentors, mentor_positions):
        # Card
        add_rect(slide, mx, my, card_w, card_h, WHITE_C, corner_radius=0.05)
        # Color bar top
        add_rect(slide, mx, my, card_w, Pt(5), mentor_color)
        # Mentor name
        add_text(slide, mx + Inches(0.15), my + Inches(0.15), card_w - Inches(0.3), Inches(0.35),
                 mentor_name, font_size=14, bold=True, color=DARK)
        add_text(slide, mx + Inches(0.15), my + Inches(0.45), card_w - Inches(0.3), Inches(0.2),
                 f"Mentor  \u2014  {len(startups)} startups assign\u00e9es", font_size=8, color=MUTED)
        # Startups list
        for j, startup in enumerate(startups):
            y_off = my + Inches(0.75) + j * Inches(0.3)
            # Bullet with startup name
            add_rect(slide, mx + Inches(0.2), y_off + Inches(0.05), Inches(0.12), Inches(0.12), mentor_color, corner_radius=0.5)
            add_text(slide, mx + Inches(0.4), y_off, card_w - Inches(0.6), Inches(0.25),
                     startup, font_size=10, color=DARK)

    # Summary KPIs
    add_text(slide, Inches(0.4), Inches(6.85), Inches(9), Inches(0.25),
             "4 mentors  |  14 startups accompagn\u00e9es  |  Ratio moyen : 3,5 startups / mentor",
             font_size=9, bold=True, color=SLATE, alignment=PP_ALIGN.CENTER)

    # Footer
    add_text(slide, Inches(0.4), Inches(7.15), Inches(6), Inches(0.2),
             "The Pulse  |  UM6P  |  thepulse.ma  |  Rapport généré le 13 avril 2026",
             font_size=8, color=MUTED)
    add_text(slide, Inches(6.5), Inches(7.15), Inches(3.2), Inches(0.2),
             "Page 8/8  |  Rapport partenaires",
             font_size=8, color=MUTED, alignment=PP_ALIGN.RIGHT)

    prs.save(OUTPUT)
    print(f"[OK] Rapport PPTX généré : {OUTPUT}")


if __name__ == "__main__":
    build_pptx()
